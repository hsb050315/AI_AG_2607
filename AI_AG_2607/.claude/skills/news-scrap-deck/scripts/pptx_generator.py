"""python-pptx 기반 심플 발표자료(16:9) 생성 템플릿."""

import copy
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
DEFAULT_MAIN_COLOR = "1F4E79"  # 남색 계열 (docx_generator.py와 동일 기본값)
FONT_NAME = "맑은 고딕"


def _hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor.from_string(hex_color.upper())


def _clear_slides(prs):
    """Presentation에 이미 존재하는 슬라이드를 모두 제거한다 (템플릿 원본 슬라이드 제거용)."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _pick_blank_layout(prs):
    """플레이스홀더가 0~1개인 레이아웃(빈 레이아웃에 가까운 것)을 찾는다. 없으면 마지막 레이아웃을 사용."""
    best = None
    best_count = None
    for layout in prs.slide_layouts:
        count = len(layout.placeholders)
        if best_count is None or count < best_count:
            best, best_count = layout, count
        if count <= 1:
            return layout
    return best if best is not None else prs.slide_layouts[-1]


def _extract_theme_accent1(prs):
    """템플릿 테마의 accent1 색상(hex)을 추출한다. 실패 시 None 반환."""
    try:
        theme_el = prs.slide_masters[0].element.getroottree().getroot()
    except Exception:
        return None
    try:
        theme_part = prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        accent1 = theme_part.element.find(".//a:accent1/a:srgbClr", ns)
        if accent1 is not None:
            return accent1.get("val")
    except Exception:
        return None
    return None


def duplicate_slide(prs, index):
    """prs 내 index번째 슬라이드를 복제해 맨 뒤에 추가한다 (이미지/서식 포함, 같은 prs 안에서만 동작).

    Canva 등에서 내보낸 pptx는 배경/그래픽이 슬라이드 레이아웃이 아니라 각 슬라이드에 직접
    그려져 있는 경우가 많다. 이런 템플릿의 비주얼을 재사용하려면 `new_presentation(template_path=...)`
    로 연 뒤 이 함수로 실제 슬라이드를 복제하고, `remove_shapes`/`set_shape_text`/`set_shape_bullets`로
    복제본의 도형만 다듬어 쓰는 편이 `template_path`의 빈 레이아웃 방식보다 훨씬 원본에 가깝다.
    """
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))

    # 관계(rId)는 대상 파트에 새로 추가되며 원본과 다른 rId를 받을 수 있으므로,
    # old rId -> new rId 매핑을 만들어 복제된 도형 XML의 r:id/r:embed 참조를 갱신한다.
    rid_map = {}
    for rid, rel in source.part.rels.items():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rid = dest.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rid] = new_rid

    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for el in dest.shapes._spTree.iter():
        for attr in list(el.attrib):
            if attr.startswith("{%s}" % r_ns):
                old_rid = el.attrib[attr]
                if old_rid in rid_map:
                    el.attrib[attr] = rid_map[old_rid]
    return dest


def delete_slide(prs, index):
    """prs에서 index번째 슬라이드를 제거한다."""
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def remove_shapes(slide, names):
    """이름(shape.name)이 names에 포함되는 도형을 슬라이드에서 제거한다."""
    for shp in list(slide.shapes):
        if shp.name in names:
            shp._element.getparent().remove(shp._element)


def set_shape_text(shape, text):
    """도형의 첫 번째 런(run) 서식은 유지한 채 텍스트만 교체하고, 나머지 문단은 제거한다."""
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    first = paras[0]
    if first.runs:
        first.runs[0].text = text
        for r in first.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        first.text = text
    for p in paras[1:]:
        p._p.getparent().remove(p._p)


def set_shape_bullets(shape, lines, bullet_prefix="•  "):
    """도형의 첫 문단 서식을 복제해 여러 줄(bullet) 텍스트를 채운다."""
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    first = paras[0]
    for p in paras[1:]:
        p._p.getparent().remove(p._p)
    set_shape_text(shape, f"{bullet_prefix}{lines[0]}" if lines else "")
    first = tf.paragraphs[0]
    for line in lines[1:]:
        new_p_el = copy.deepcopy(first._p)
        first._p.getparent().append(new_p_el)
        new_p = tf.paragraphs[-1]
        if new_p.runs:
            new_p.runs[0].text = f"{bullet_prefix}{line}"
            for r in new_p.runs[1:]:
                r._r.getparent().remove(r._r)
        else:
            new_p.text = f"{bullet_prefix}{line}"


def new_presentation(main_color=None, template_path=None):
    """새 Presentation 객체 생성. main_color(hex)로 전체 강조색을 지정한다.

    template_path가 주어지면 해당 pptx를 베이스로 열어 슬라이드 크기/테마를 유지하고
    기존 슬라이드는 모두 제거한 뒤 빈 캔버스로 반환한다. 미지정 시 기존과 동일하게
    새 16:9 Presentation을 생성한다.
    """
    if template_path:
        prs = Presentation(template_path)
        _clear_slides(prs)
        prs._blank_layout = _pick_blank_layout(prs)
        if main_color:
            prs.main_color = _hex_to_rgb(main_color)
        else:
            accent1 = _extract_theme_accent1(prs)
            prs.main_color = _hex_to_rgb(accent1) if accent1 else _hex_to_rgb(DEFAULT_MAIN_COLOR)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        prs._blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃
        prs.main_color = _hex_to_rgb(main_color or DEFAULT_MAIN_COLOR)
    return prs


def _blank_slide(prs):
    layout = getattr(prs, "_blank_layout", None) or prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _set_font(run, size=None, bold=False, color=None, italic=False):
    run.font.name = FONT_NAME
    if size:
        run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def add_title_slide(prs, title, subtitle=None):
    """표지 슬라이드: 중앙 정렬 제목 + 부제."""
    slide = _blank_slide(prs)

    box = slide.shapes.add_textbox(Inches(1), Inches(2.7), prs.slide_width - Inches(2), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    _set_font(run, size=Pt(44), bold=True, color=prs.main_color)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), prs.slide_width - Inches(2), Inches(1))
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = subtitle
        _set_font(srun, size=Pt(20), color=RGBColor(0x66, 0x66, 0x66))

    return slide


def _add_header(slide, prs, title):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), prs.slide_width - Inches(1.2), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_font(run, size=Pt(28), bold=True, color=prs.main_color)

    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.15), prs.slide_width - Inches(1.2), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = prs.main_color
    line.line.fill.background()
    return box


def add_bullet_slide(prs, title, bullets, source=None):
    """제목 + 글머리 기호 목록 슬라이드."""
    slide = _blank_slide(prs)
    _add_header(slide, prs, title)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), prs.slide_width - Inches(1.8), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = f"•  {item}"
        _set_font(run, size=Pt(20))
        p.space_after = Pt(14)

    if source:
        src_box = slide.shapes.add_textbox(Inches(0.9), Inches(6.9), prs.slide_width - Inches(1.8), Inches(0.5))
        stf = src_box.text_frame
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = source
        _set_font(srun, size=Pt(12), italic=True, color=RGBColor(0x88, 0x88, 0x88))

    return slide


def add_image_bullet_slide(prs, title, bullets, image_path, source=None):
    """왼쪽 텍스트(글머리 기호) + 오른쪽 스크린샷 썸네일 슬라이드."""
    slide = _blank_slide(prs)
    _add_header(slide, prs, title)

    text_width = Emu(int(prs.slide_width * 0.49))
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), text_width, Inches(5.0))
    tf = text_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"•  {item}"
        _set_font(run, size=Pt(18))
        p.space_after = Pt(12)

    if image_path and os.path.exists(image_path):
        pic_width = Emu(int(prs.slide_width * 0.37))
        pic_left = prs.slide_width - pic_width - Inches(0.7)
        pic_top = Inches(1.6)
        slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width)

    if source:
        src_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), prs.slide_width - Inches(1.4), Inches(0.5))
        stf = src_box.text_frame
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = source
        _set_font(srun, size=Pt(12), italic=True, color=RGBColor(0x88, 0x88, 0x88))

    return slide


def add_section_slide(prs, title):
    """구간을 나누는 심플 섹션 구분 슬라이드 (배경색 + 중앙 제목)."""
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = prs.main_color
    bg.line.fill.background()
    bg.shadow.inherit = False

    box = slide.shapes.add_textbox(Inches(1), Inches(3.2), prs.slide_width - Inches(2), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    _set_font(run, size=Pt(36), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    return slide


def add_closing_slide(prs, message="감사합니다"):
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(1), Inches(3.2), prs.slide_width - Inches(2), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = message
    _set_font(run, size=Pt(32), bold=True, color=prs.main_color)
    return slide


def save(prs, path):
    parent = os.path.dirname(path)
    if parent:
        os.makedirs(parent, exist_ok=True)
    prs.save(path)
    print(f"저장 완료: {path}")


def build_sample_presentation(output_path="sample_output.pptx", main_color=None):
    """템플릿 기능을 확인할 수 있는 예시 발표자료 생성."""
    prs = new_presentation(main_color=main_color)

    add_title_slide(prs, "OO 이슈 뉴스 스크랩 브리핑", subtitle="조사일자: 2026년 8월 11일")
    add_section_slide(prs, "1. 기사별 요약")
    add_bullet_slide(
        prs,
        "기사 제목 예시",
        ["핵심 요약 1", "핵심 요약 2", "핵심 요약 3"],
        source="출처: OO뉴스 · 2026-08-11",
    )
    add_section_slide(prs, "2. 종합 인사이트")
    add_bullet_slide(prs, "종합 결론", ["트렌드 1", "트렌드 2", "시사점"])
    add_closing_slide(prs)

    save(prs, output_path)


if __name__ == "__main__":
    build_sample_presentation(main_color="2E86AB")
